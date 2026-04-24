"""
Gera a apresentação PowerPoint do JogoScore para a Direção.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Paleta de cores ──────────────────────────────────────────
ROXO_ESC   = RGBColor(0x4C, 0x1D, 0x95)   # roxo escuro
ROXO       = RGBColor(0x76, 0x4B, 0xA2)   # roxo médio
AZUL       = RGBColor(0x16, 0x7E, 0xEA)   # azul
VERDE      = RGBColor(0x05, 0x96, 0x69)   # verde IDOntime
VERDE_CLR  = RGBColor(0xD1, 0xFA, 0xE5)   # verde claro
AMARELO    = RGBColor(0xF5, 0x9E, 0x0B)   # amarelo bónus
BRANCO     = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLR  = RGBColor(0xF3, 0xF4, 0xF6)
CINZA_ESC  = RGBColor(0x1F, 0x2A, 0x44)
CINZA_MED  = RGBColor(0x6B, 0x72, 0x80)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completamente em branco


# ── Helpers ──────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=BRANCO,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullet_box(slide, items, x, y, w, h,
                   font_size=16, color=CINZA_ESC,
                   bullet="•", line_spacing=1.2):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
    return txb


def gradient_bg(slide, color1, color2):
    """Fundo com gradiente simulado via dois rectângulos."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color1)
    # rectângulo diagonal semi-transparente com a 2ª cor
    shape = slide.shapes.add_shape(1,
        Inches(7), 0, Inches(6.33), SLIDE_H)
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = color2
    # transparência via XML
    sp_pr = shape.fill._xPr
    solid = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solid is not None:
        srgb = solid.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_el = etree.SubElement(
                srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_el.set('val', '40000')   # 40% opacidade


def slide_titulo():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, ROXO_ESC)
    gradient_bg(slide, ROXO_ESC, ROXO)

    # Barra lateral esquerda decorativa
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, AMARELO)

    # Emoji / ícone
    add_text(slide, "🎮", Inches(0.6), Inches(0.8), Inches(2), Inches(1.5),
             font_size=52, align=PP_ALIGN.LEFT)

    # Título principal
    add_text(slide, "JogoScore", Inches(0.6), Inches(1.9), Inches(9), Inches(1.4),
             font_size=60, bold=True, color=BRANCO)

    # Subtítulo
    add_text(slide, "Sistema de Gamificação de Colaboradores — SOCEM",
             Inches(0.6), Inches(3.05), Inches(9), Inches(0.7),
             font_size=22, color=RGBColor(0xC4, 0xB5, 0xFD))

    # Linha separadora
    add_rect(slide, Inches(0.6), Inches(3.85), Inches(5), Inches(0.04), AMARELO)

    # Info rodapé
    add_text(slide, "Apresentação à Direção  ·  2026",
             Inches(0.6), Inches(4.1), Inches(8), Inches(0.5),
             font_size=14, color=RGBColor(0xA7, 0x8B, 0xFA), italic=True)

    # Tags SINEX + IDOntime
    add_rect(slide, Inches(0.6), Inches(5.2), Inches(1.5), Inches(0.45), AZUL)
    add_text(slide, "SINEX", Inches(0.6), Inches(5.2), Inches(1.5), Inches(0.45),
             font_size=13, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(2.3), Inches(5.2), Inches(1.8), Inches(0.45), VERDE)
    add_text(slide, "IDOntime", Inches(2.3), Inches(5.2), Inches(1.8), Inches(0.45),
             font_size=13, bold=True, align=PP_ALIGN.CENTER)


def slide_o_que_e():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_CLR)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), ROXO_ESC)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, ROXO)

    add_text(slide, "O que é o JogoScore?",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    # Caixas de conceito
    conceitos = [
        ("🎯", "Motivação",
         "Transforma o cumprimento de regras e objetivos do dia-a-dia em pontos e recompensas visíveis."),
        ("🔗", "Integração",
         "Liga-se automaticamente ao SINEX (assiduidade) e ao IDOntime (picagens) sem intervenção manual."),
        ("🏆", "Competição Saudável",
         "Ranking individual e por departamento visible a todos, promovendo o espírito de equipa."),
        ("🛍️", "Loja de Recompensas",
         "Os pontos ganhos podem ser trocados por prémios reais definidos pela empresa."),
    ]

    col_w = Inches(2.9)
    gap   = Inches(0.22)
    for i, (icon, titulo, desc) in enumerate(conceitos):
        x = Inches(0.35) + i * (col_w + gap)
        add_rect(slide, x, Inches(1.75), col_w, Inches(4.8), BRANCO)
        # barra topo colorida
        cores = [ROXO, AZUL, VERDE, AMARELO]
        add_rect(slide, x, Inches(1.75), col_w, Inches(0.12), cores[i])
        add_text(slide, icon, x, Inches(1.95), col_w, Inches(0.9),
                 font_size=36, align=PP_ALIGN.CENTER)
        add_text(slide, titulo, x, Inches(2.85), col_w, Inches(0.55),
                 font_size=17, bold=True, color=CINZA_ESC, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.15), Inches(3.5),
                 col_w - Inches(0.3), Inches(2.8),
                 font_size=13, color=CINZA_MED, align=PP_ALIGN.LEFT)


def slide_sistemas():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_CLR)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), ROXO_ESC)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, ROXO)

    add_text(slide, "Sistemas Integrados",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    # SINEX — lado esquerdo
    add_rect(slide, Inches(0.35), Inches(1.7), Inches(5.9), Inches(5.3), BRANCO)
    add_rect(slide, Inches(0.35), Inches(1.7), Inches(5.9), Inches(0.55), AZUL)
    add_text(slide, "📊  SINEX — Assiduidade",
             Inches(0.5), Inches(1.72), Inches(5.7), Inches(0.5),
             font_size=16, bold=True, color=BRANCO)

    sinex_items = [
        "Registo diário de presenças via sistema SINEX",
        "Processamento automático todas as manhãs",
        "Analisa se o colaborador esteve presente e cumpriu o turno",
        "+1 ponto por dia de presença válida",
        "-1 ponto por falta ou ausência injustificada",
        "Bónus de +1 pela semana completa (seg → sex) sem faltas",
        "Exclui feriados nacionais e municipal (Alcobaça 20 Ago)",
    ]
    add_bullet_box(slide, sinex_items,
                   Inches(0.5), Inches(2.4), Inches(5.6), Inches(4.3),
                   font_size=13.5, color=CINZA_ESC)

    # IDOntime — lado direito
    add_rect(slide, Inches(6.55), Inches(1.7), Inches(6.4), Inches(5.3), BRANCO)
    add_rect(slide, Inches(6.55), Inches(1.7), Inches(6.4), Inches(0.55), VERDE)
    add_text(slide, "🕐  IDOntime — Picagens de Ponto",
             Inches(6.7), Inches(1.72), Inches(6.2), Inches(0.5),
             font_size=16, bold=True, color=BRANCO)

    idon_items = [
        "Leitura das picagens de entrada/saída do IDOnics",
        "Processamento automático todas as manhãs",
        "Verifica se o colaborador atingiu o objetivo de horas",
        "+1 ponto se as horas efectivas >= objetivo do plano",
        "-1 ponto se tem picagens mas não atinge o objetivo",
        "Regra dos 30 min: tolerância por entrada tardia/saída cedo",
        "Turnos nocturnos: combina picagens do dia D e D+1",
        "Bónus de +1 pela semana completa (seg → sex) sem falhas",
    ]
    add_bullet_box(slide, idon_items,
                   Inches(6.7), Inches(2.4), Inches(6.1), Inches(4.3),
                   font_size=13.5, color=CINZA_ESC)


def slide_pontos_sinex():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xEE, 0xF2, 0xFF))
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), AZUL)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, ROXO)

    add_text(slide, "📊  Lógica de Pontos — SINEX",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    # Fluxo em caixas
    steps = [
        (AZUL,    "1. Fonte de Dados",
                  "O job lê o SINEX diariamente.\nVerifica presenças e ausências de cada colaborador."),
        (ROXO,    "2. Cálculo do Ponto",
                  "+1 → Presente e com registo válido\n−1 → Ausência ou falta\n 0 → Feriado ou dia não útil (ignorado)"),
        (AMARELO, "3. Bónus Semanal",
                  "Na segunda-feira verifica a semana anterior.\nSe teve +1 de seg a sex → recebe +1 extra."),
        (VERDE,   "4. Resultado Final",
                  "Pontos acumulados no ranking individual\ne ranking por departamento."),
    ]

    bx_w = Inches(2.8)
    bx_h = Inches(3.8)
    gap  = Inches(0.32)
    start_x = Inches(0.35)

    for i, (cor, titulo, desc) in enumerate(steps):
        x = start_x + i * (bx_w + gap)
        add_rect(slide, x, Inches(1.75), bx_w, bx_h, BRANCO)
        add_rect(slide, x, Inches(1.75), bx_w, Inches(0.55), cor)
        add_text(slide, titulo, x + Inches(0.15), Inches(1.78),
                 bx_w - Inches(0.3), Inches(0.5),
                 font_size=14, bold=True, color=BRANCO)
        add_text(slide, desc, x + Inches(0.15), Inches(2.45),
                 bx_w - Inches(0.3), Inches(2.9),
                 font_size=13, color=CINZA_ESC)
        # seta entre caixas
        if i < len(steps) - 1:
            ax = x + bx_w + Inches(0.05)
            add_text(slide, "→", ax, Inches(2.9), Inches(0.25), Inches(0.5),
                     font_size=22, bold=True, color=ROXO, align=PP_ALIGN.CENTER)

    # Nota de rodapé
    add_rect(slide, Inches(0.35), Inches(5.8), Inches(12.6), Inches(0.9),
             RGBColor(0xE0, 0xE7, 0xFF))
    add_text(slide,
             "ℹ️  Anti-duplicação garantida: cada dia só pode ter um registo por colaborador. "
             "O sistema faz UPSERT — se re-processar, atualiza sem duplicar.",
             Inches(0.55), Inches(5.82), Inches(12.2), Inches(0.85),
             font_size=12, color=AZUL)


def slide_pontos_idontime():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xEC, 0xFD, 0xF5))
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), VERDE)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, RGBColor(0x06, 0x4E, 0x3B))

    add_text(slide, "🕐  Lógica de Pontos — IDOntime",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    # Tabela de regras
    add_rect(slide, Inches(0.35), Inches(1.7), Inches(6.1), Inches(4.5), BRANCO)
    add_rect(slide, Inches(0.35), Inches(1.7), Inches(6.1), Inches(0.5), VERDE)
    add_text(slide, "Regras de Pontuação Diária",
             Inches(0.5), Inches(1.72), Inches(5.9), Inches(0.46),
             font_size=15, bold=True, color=BRANCO)

    regras = [
        ("✅  +1 ponto",  "Horas efectivas ≥ objetivo do plano de trabalho"),
        ("❌  −1 ponto",  "Tem picagens mas não atinge o mínimo de horas"),
        ("⏭   0 pontos",  "Sem picagens (folga ou ausência) — ignorado"),
        ("⭐  +1 bónus",  "Semana completa sem falhas (seg → sex) — segunda-feira"),
        ("✏️  Manual",    "Admin pode atribuir pontos manuais com motivo"),
    ]

    for i, (pts, desc) in enumerate(regras):
        y = Inches(2.35) + i * Inches(0.72)
        bg = CINZA_CLR if i % 2 == 0 else BRANCO
        add_rect(slide, Inches(0.35), y, Inches(6.1), Inches(0.68), bg)
        add_text(slide, pts,  Inches(0.5),  y + Inches(0.08),
                 Inches(1.6), Inches(0.55), font_size=13, bold=True, color=CINZA_ESC)
        add_text(slide, desc, Inches(2.05), y + Inches(0.08),
                 Inches(4.2), Inches(0.55), font_size=12.5, color=CINZA_MED)

    # Regra dos 30 min
    add_rect(slide, Inches(6.75), Inches(1.7), Inches(6.15), Inches(2.15), BRANCO)
    add_rect(slide, Inches(6.75), Inches(1.7), Inches(6.15), Inches(0.5),
             RGBColor(0x06, 0x4E, 0x3B))
    add_text(slide, "⏱  Regra dos 30 Minutos",
             Inches(6.9), Inches(1.72), Inches(5.95), Inches(0.46),
             font_size=15, bold=True, color=BRANCO)
    r30 = [
        "Entrada tardia → acrescenta 30 min por fracção",
        "Saída antecipada → desconta 30 min por fracção",
        "Regresso tardio da pausa → acrescenta 30 min",
        "Entrada antecipada → sem bónus (conta hora marcada)",
    ]
    add_bullet_box(slide, r30, Inches(6.9), Inches(2.3),
                   Inches(5.8), Inches(1.4), font_size=12.5, color=CINZA_ESC)

    # Turnos nocturnos
    add_rect(slide, Inches(6.75), Inches(4.0), Inches(6.15), Inches(2.2), BRANCO)
    add_rect(slide, Inches(6.75), Inches(4.0), Inches(6.15), Inches(0.5),
             RGBColor(0x1E, 0x40, 0xAF))
    add_text(slide, "🌙  Turnos Nocturnos",
             Inches(6.9), Inches(4.02), Inches(5.95), Inches(0.46),
             font_size=15, bold=True, color=BRANCO)
    noct = [
        "Entrada ≥ 20:00 → turno nocturno",
        "Combina picagens do dia D com as do dia D+1",
        "Garante contagem correcta de horas contínuas",
    ]
    add_bullet_box(slide, noct, Inches(6.9), Inches(4.6),
                   Inches(5.8), Inches(1.45), font_size=12.5, color=CINZA_ESC)


def slide_bonus():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xFF, 0xFB, 0xEB))
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), AMARELO)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, RGBColor(0x92, 0x40, 0x0E))

    add_text(slide, "⭐  Bónus de Semana Completa",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    # Explicação central
    add_rect(slide, Inches(0.35), Inches(1.7), Inches(12.6), Inches(1.2),
             RGBColor(0xFF, 0xF3, 0xC4))
    add_text(slide,
             "Todos os colaboradores (SINEX e IDOntime) que cumprirem todos os dias úteis "
             "de segunda a sexta sem nenhuma penalização recebem automaticamente +1 ponto bónus na segunda-feira seguinte.",
             Inches(0.55), Inches(1.78), Inches(12.2), Inches(1.1),
             font_size=15, color=RGBColor(0x78, 0x35, 0x00))

    # Calendário visual
    dias = ["Segunda", "Terça", "Quarta", "Quinta", "Sexta"]
    cores_ok   = VERDE
    w_dia = Inches(1.95)
    gap   = Inches(0.22)
    sx    = Inches(0.7)

    for i, dia in enumerate(dias):
        x = sx + i * (w_dia + gap)
        add_rect(slide, x, Inches(3.1), w_dia, Inches(1.8), VERDE)
        add_text(slide, dia,  x, Inches(3.15), w_dia, Inches(0.5),
                 font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        add_text(slide, "+1", x, Inches(3.65), w_dia, Inches(0.8),
                 font_size=32, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    # Seta + bónus
    add_text(slide, "→", Inches(11.0), Inches(3.55), Inches(0.5), Inches(0.8),
             font_size=28, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(11.6), Inches(3.1), Inches(1.35), Inches(1.8), AMARELO)
    add_text(slide, "⭐", Inches(11.6), Inches(3.15), Inches(1.35), Inches(0.6),
             font_size=24, align=PP_ALIGN.CENTER)
    add_text(slide, "Bónus\n+1", Inches(11.6), Inches(3.7), Inches(1.35), Inches(0.9),
             font_size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    # Regras
    add_rect(slide, Inches(0.35), Inches(5.15), Inches(12.6), Inches(1.9), BRANCO)
    notas = [
        "Processado automaticamente na segunda-feira pelo job diário (SINEX e IDOntime)",
        "Se a semana tiver feriado(s), tem menos dias úteis — o bónus não é atribuído nessa semana",
        "Anti-duplicação: o bónus só é dado uma vez por semana por colaborador",
        "Também é aplicado ao reprocessar intervalos históricos (recálculo retroativo)",
    ]
    add_bullet_box(slide, notas, Inches(0.55), Inches(5.2),
                   Inches(12.2), Inches(1.8), font_size=13, color=CINZA_ESC)


def slide_funcionalidades():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_CLR)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), ROXO_ESC)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, AMARELO)

    add_text(slide, "Funcionalidades da Plataforma",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    funcs = [
        ("🏠", "Dashboard",
         "Vista geral com ranking global,\nleaderboard por pontos e\nacesso rápido a todas as áreas."),
        ("🏢", "Departamentos",
         "Ranking de departamentos com\npontos totais, média por pessoa\ne top 3 destacado."),
        ("👤", "Perfil Colaborador",
         "Histórico de pontos individual,\nranking no departamento e\nna empresa."),
        ("🛍️", "Loja de Prémios",
         "Catálogo de recompensas.\nColaborador gasta pontos\npara resgatar prémios."),
        ("📊", "Relatório Mensal",
         "Análise por mês: melhor\ncolaborador, mais consistente\ne stats por departamento."),
        ("⚙️", "Área de Admin",
         "Gestão de contas, atribuição\nmanual de pontos, reprocessamento\ne sincronização."),
    ]

    col_w = Inches(3.85)
    row_h = Inches(2.4)
    gap_x = Inches(0.25)
    gap_y = Inches(0.2)
    sx = Inches(0.35)
    sy = Inches(1.7)

    for i, (icon, titulo, desc) in enumerate(funcs):
        col = i % 3
        row = i // 3
        x = sx + col * (col_w + gap_x)
        y = sy + row * (row_h + gap_y)
        add_rect(slide, x, y, col_w, row_h, BRANCO)
        # barra lateral colorida
        cores = [ROXO, AZUL, VERDE, AMARELO, RGBColor(0xEC,0x48,0x99), RGBColor(0x06,0x95,0x69)]
        add_rect(slide, x, y, Inches(0.1), row_h, cores[i])
        add_text(slide, icon,   x + Inches(0.25), y + Inches(0.2),
                 Inches(0.7), Inches(0.7), font_size=26)
        add_text(slide, titulo, x + Inches(0.25), y + Inches(0.85),
                 col_w - Inches(0.4), Inches(0.5),
                 font_size=15, bold=True, color=CINZA_ESC)
        add_text(slide, desc,   x + Inches(0.25), y + Inches(1.3),
                 col_w - Inches(0.4), Inches(0.95),
                 font_size=12, color=CINZA_MED)


def slide_admin():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_CLR)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), ROXO_ESC)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, AZUL)

    add_text(slide, "⚙️  Área de Administração",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    grupos = [
        ("👥  Gestão de Utilizadores", ROXO, [
            "Colaboradores SINEX — criar, listar, ver histórico",
            "Colaboradores IDOntime — associar a departamentos",
            "Contas de Login — criar acesso individual à plataforma",
        ]),
        ("💰  Atribuição de Pontos", AZUL, [
            "Atribuição manual de pontos a qualquer colaborador",
            "Suporte a SINEX e IDOntime em separado",
            "Motivo obrigatório, anti-duplicação por UUID",
        ]),
        ("🔄  Reprocessamento", VERDE, [
            "Reprocessar SINEX para qualquer intervalo de datas",
            "Reprocessar IDOntime com recálculo de picagens",
            "Bónus semanal retroativo incluído no reprocessamento",
        ]),
        ("📊  Relatórios", AMARELO, [
            "Relatório mensal com top colaboradores e destaques",
            "Stats por departamento com total e média de pontos",
            "Exportação para Excel",
        ]),
    ]

    col_w = Inches(5.85)
    col_h = Inches(2.55)
    gap   = Inches(0.25)
    sx    = Inches(0.35)

    for i, (titulo, cor, items) in enumerate(grupos):
        col = i % 2
        row = i // 2
        x = sx + col * (col_w + gap)
        y = Inches(1.75) + row * (col_h + gap)
        add_rect(slide, x, y, col_w, col_h, BRANCO)
        add_rect(slide, x, y, col_w, Inches(0.5), cor)
        add_text(slide, titulo, x + Inches(0.15), y + Inches(0.03),
                 col_w - Inches(0.3), Inches(0.46),
                 font_size=14, bold=True, color=BRANCO)
        add_bullet_box(slide, items, x + Inches(0.15), y + Inches(0.6),
                       col_w - Inches(0.3), col_h - Inches(0.75),
                       font_size=12.5, color=CINZA_ESC)


def slide_arquitetura():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CINZA_CLR)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), ROXO_ESC)
    add_rect(slide, 0, 0, Inches(0.1), SLIDE_H, VERDE)

    add_text(slide, "🏗️  Arquitectura Técnica",
             Inches(0.4), Inches(0.28), Inches(10), Inches(1),
             font_size=30, bold=True, color=BRANCO)

    # Diagrama simplificado
    blocos = [
        (Inches(0.5),  Inches(2.2), Inches(2.4), Inches(1.1), AZUL,   "SINEX\n(Base de dados)"),
        (Inches(0.5),  Inches(3.8), Inches(2.4), Inches(1.1), VERDE,  "IDOnics\n(Base de dados)"),
        (Inches(4.0),  Inches(2.0), Inches(2.8), Inches(1.4), ROXO,   "sinex_job.py\nJob diário automático"),
        (Inches(4.0),  Inches(3.7), Inches(2.8), Inches(1.4), VERDE,  "idontime_job.py\nJob diário automático"),
        (Inches(7.5),  Inches(2.7), Inches(2.6), Inches(1.2), RGBColor(0x1D,0x40,0xAF), "jogo_score\njogo_socem\njogo_idonics"),
        (Inches(10.5), Inches(2.7), Inches(2.3), Inches(1.2), ROXO_ESC, "JogoScore\nFastAPI + Web"),
    ]

    for x, y, w, h, cor, texto in blocos:
        add_rect(slide, x, y, w, h, cor)
        add_text(slide, texto, x + Inches(0.1), y + Inches(0.1),
                 w - Inches(0.2), h - Inches(0.2),
                 font_size=12.5, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    # Setas (texto)
    setas = [
        (Inches(2.95), Inches(2.6),  "→"),
        (Inches(2.95), Inches(4.25), "→"),
        (Inches(6.85), Inches(3.1),  "→"),
        (Inches(10.1), Inches(3.1),  "→"),
    ]
    for x, y, s in setas:
        add_text(slide, s, x, y, Inches(0.6), Inches(0.55),
                 font_size=22, bold=True, color=ROXO, align=PP_ALIGN.CENTER)

    # Task Scheduler
    add_rect(slide, Inches(3.8), Inches(5.4), Inches(3.2), Inches(0.8),
             RGBColor(0xF3,0xF4,0xF6))
    add_rect(slide, Inches(3.8), Inches(5.4), Inches(3.2), Inches(0.08), AMARELO)
    add_text(slide, "🗓  Windows Task Scheduler\nCorre automaticamente todos os dias de manhã",
             Inches(3.95), Inches(5.45), Inches(3.0), Inches(0.75),
             font_size=11, color=CINZA_ESC)

    # Tecnologias
    add_rect(slide, Inches(0.35), Inches(5.55), Inches(3.0), Inches(1.5), BRANCO)
    add_text(slide, "🛠  Tecnologias",
             Inches(0.5), Inches(5.58), Inches(2.8), Inches(0.4),
             font_size=13, bold=True, color=CINZA_ESC)
    techs = ["Python · FastAPI · Jinja2", "SQL Server · pyodbc · pandas",
             "Tailwind CSS · HTML5"]
    add_bullet_box(slide, techs, Inches(0.5), Inches(6.0),
                   Inches(2.8), Inches(0.95), font_size=11.5, color=CINZA_MED, bullet="▸")

    # Browsers
    add_rect(slide, Inches(10.3), Inches(4.15), Inches(2.7), Inches(0.9), BRANCO)
    add_text(slide, "🌐  Acesso via browser\nhttp://servidor:8005",
             Inches(10.45), Inches(4.18), Inches(2.5), Inches(0.85),
             font_size=12, color=CINZA_ESC)


def slide_final():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, ROXO_ESC)
    gradient_bg(slide, ROXO_ESC, ROXO)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, AMARELO)

    add_text(slide, "🎮", Inches(0.6), Inches(1.2), Inches(2), Inches(1.2),
             font_size=52, align=PP_ALIGN.LEFT)

    add_text(slide, "JogoScore está pronto.",
             Inches(0.6), Inches(2.2), Inches(12), Inches(1.1),
             font_size=44, bold=True, color=BRANCO)

    add_text(slide,
             "Um sistema que transforma o dia-a-dia dos colaboradores\n"
             "num jogo justo, transparente e motivador.",
             Inches(0.6), Inches(3.35), Inches(11), Inches(1.1),
             font_size=20, color=RGBColor(0xC4, 0xB5, 0xFD))

    add_rect(slide, Inches(0.6), Inches(4.55), Inches(4.5), Inches(0.04), AMARELO)

    pontos_finais = [
        "✅  Integração automática com SINEX e IDOntime",
        "✅  Bónus semanal para incentivar consistência",
        "✅  Loja de recompensas com pontos reais",
        "✅  Relatórios e administração completos",
    ]
    add_bullet_box(slide, pontos_finais,
                   Inches(0.6), Inches(4.75), Inches(9), Inches(2.2),
                   font_size=16, color=RGBColor(0xE0, 0xE7, 0xFF), bullet="")


# ── Gerar todos os slides ─────────────────────────────────────

slide_titulo()
slide_o_que_e()
slide_sistemas()
slide_pontos_sinex()
slide_pontos_idontime()
slide_bonus()
slide_funcionalidades()
slide_admin()
slide_arquitetura()
slide_final()

output = "JogoScore_Apresentacao.pptx"
prs.save(output)
print(f"Apresentacao gerada: {output}")
